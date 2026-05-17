#!/usr/bin/env python3
"""Render a PPTX file from a JSON payload.

Reads JSON from stdin, writes binary .pptx to stdout, errors to stderr.

Payload shape:
{
  "theme": {
    "aspect": [10.0, 7.5],
    "bg": "#000000",
    "text": "#FFFFFF",
    "font": "Calibri",
    "size": 42,
    "h_align": "center",
    "v_align": "middle",
    "margins": 0.05,
    "bold_section_types": ["chorus"]
  },
  "slides": [
    {"type": "chorus", "lines": ["AMÉM", "AMÉM"]},
    {"type": "verse",  "lines": ["Vem e eu mostrarei", "..."]},
    {"type": "blank",  "lines": []}
  ]
}
"""
import io
import json
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


H_ALIGN = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}

V_ALIGN = {
    "top": MSO_ANCHOR.TOP,
    "middle": MSO_ANCHOR.MIDDLE,
    "bottom": MSO_ANCHOR.BOTTOM,
}


def hex_to_rgb(value):
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def add_blank_slide(prs, bg_color):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    return slide


def render_slide(prs, theme, entry, bg_color, text_color, h_align, v_align):
    slide = add_blank_slide(prs, bg_color)

    lines = entry.get("lines") or []
    if not lines:
        return  # Blank slide already has the background fill

    margin = float(theme.get("margins", 0.05))
    aspect_w, aspect_h = theme.get("aspect", [10.0, 7.5])
    left = Inches(aspect_w * margin)
    top = Inches(aspect_h * margin)
    width = Inches(aspect_w * (1 - 2 * margin))
    height = Inches(aspect_h * (1 - 2 * margin))

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = v_align

    bold = entry.get("type") in theme.get("bold_section_types", [])
    font_name = theme.get("font", "Calibri")
    font_size = Pt(int(theme.get("size", 42)))

    for index, line in enumerate(lines):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.alignment = h_align
        run = paragraph.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = text_color


def render(payload):
    theme = payload.get("theme", {})
    slides = payload.get("slides", [])

    prs = Presentation()
    aspect = theme.get("aspect", [10.0, 7.5])
    prs.slide_width = Inches(aspect[0])
    prs.slide_height = Inches(aspect[1])

    bg_color = hex_to_rgb(theme.get("bg", "#000000"))
    text_color = hex_to_rgb(theme.get("text", "#FFFFFF"))
    h_align = H_ALIGN.get(theme.get("h_align", "center"), PP_ALIGN.CENTER)
    v_align = V_ALIGN.get(theme.get("v_align", "middle"), MSO_ANCHOR.MIDDLE)

    if not slides:
        # Still emit a blank deck rather than a zero-slide presentation.
        add_blank_slide(prs, bg_color)
    else:
        for entry in slides:
            render_slide(prs, theme, entry, bg_color, text_color, h_align, v_align)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def main():
    try:
        payload = json.load(sys.stdin)
    except json.JSONDecodeError as exc:
        sys.stderr.write(f"Invalid JSON payload: {exc}\n")
        return 1

    try:
        binary = render(payload)
    except Exception as exc:  # noqa: BLE001 - boundary script
        sys.stderr.write(f"Render failed: {exc}\n")
        return 1

    sys.stdout.buffer.write(binary)
    return 0


if __name__ == "__main__":
    sys.exit(main())
